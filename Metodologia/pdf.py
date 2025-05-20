from pptx import Presentation
from pptx.util import Inches

# Criando a apresentação
prs = Presentation()

# Dados dos slides com frases-chave
slides_data = [
    ("Introdução", [
        "Conceitos internos do computador",
        "Cache, Buffer, Spooling, Kernel, Processos, Memória, Arquivos"
    ]),
    ("O Kernel", [
        "Núcleo do sistema",
        "Controla tudo",
        "Gerencia recursos"
    ]),
    ("Processos", [
        "Programa em execução",
        "Estados: pronto, executando, bloqueado",
        "Kernel gerencia tudo"
    ]),
    ("Memória e Arquivos", [
        "Memória: rápida, temporária",
        "Arquivos: permanentes",
        "Kernel controla acesso"
    ]),
    ("Cache e Buffer", [
        "Cache: acesso rápido",
        "Buffer: armazenamento temporário",
        "Otimizam o desempenho"
    ]),
    ("Spooling", [
        "Fila de espera",
        "Impressão e discos",
        "Libera o usuário"
    ]),
    ("Ligando Tudo", [
        "Tudo interligado",
        "Kernel orquestra",
        "Cache, memória, arquivos, processos"
    ]),
    ("Conclusão", [
        "Componentes integrados",
        "Desempenho e eficiência",
        "Compreensão ajuda na solução de problemas"
    ]),
    ("Referências", [
        "Wikipedia, Hardware.com.br",
        "Livros de Silberschatz e Patterson",
        "Fontes confiáveis para aprofundar"
    ])
]

# Adicionando slides
for title, bullet_points in slides_data:
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    slide.shapes.title.text = title
    content = slide.shapes.placeholders[1].text = '\n'.join(bullet_points)

# Salvando a apresentação
prs.save("Apresentacao_Computador_Interno.pptx")
